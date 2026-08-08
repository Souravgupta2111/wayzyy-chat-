#!/usr/bin/env python3
"""Layout validator for Wayzyy_Moderation_Proposal_Deck.pptx.

Checks three failure modes that are invisible until the deck is opened on a
projector:

  1. bounds     - any shape extending past the slide edge
  2. overflow   - text estimated to need more height than its frame provides
  3. collision  - a text frame sitting on top of a chart or table frame

Text height is estimated, not rendered: average glyph width is taken as
0.50 * pt_size / 72 inches and line height as 1.20 * pt_size / 72. A frame is
flagged only past a 1.12x tolerance so ordinary estimation error stays quiet.
"""
import sys
from pptx import Presentation
from pptx.util import Emu

DECK = sys.argv[1] if len(sys.argv) > 1 else "/Users/apple/Desktop/Wayzyy_Moderation_Proposal_Deck.pptx"
CHAR_W_RATIO = 0.50
LINE_H_RATIO = 1.20
OVERFLOW_TOL = 1.12
EPS = 0.02


def inches(v):
    return None if v is None else Emu(v).inches


def est_height(tf, width_in):
    """Estimated rendered height, in inches, of a text frame's content."""
    total = 0.0
    for para in tf.paragraphs:
        runs = [r for r in para.runs if r.text]
        if not runs:
            total += 12.0 * LINE_H_RATIO / 72
            continue
        size = max((r.font.size.pt if r.font.size else 12.0) for r in runs)
        text = "".join(r.text for r in runs)
        char_w = CHAR_W_RATIO * size / 72
        usable = max(width_in - 0.10, 0.20)
        chars_per_line = max(int(usable / char_w), 1)
        lines = 0
        for seg in text.split("\n"):
            lines += max(1, -(-len(seg) // chars_per_line))
        spc = para.line_spacing if isinstance(para.line_spacing, float) else 1.0
        total += lines * size * LINE_H_RATIO * spc / 72
    return total


def rect(sh):
    x, y = inches(sh.left), inches(sh.top)
    w, h = inches(sh.width), inches(sh.height)
    if None in (x, y, w, h):
        return None
    return x, y, w, h


def overlaps(a, b):
    ax, ay, aw, ah = a
    bx, by, bw, bh = b
    return (ax < bx + bw - EPS and bx < ax + aw - EPS
            and ay < by + bh - EPS and by < ay + ah - EPS)


def main():
    prs = Presentation(DECK)
    SW, SH = prs.slide_width.inches, prs.slide_height.inches
    issues = []
    n_charts = n_tables = 0

    for i, slide in enumerate(prs.slides, 1):
        blocks = []
        texts = []
        for sh in slide.shapes:
            r = rect(sh)
            if r is None:
                continue
            x, y, w, h = r
            name = sh.shape_type
            if x < -EPS or y < -EPS or x + w > SW + EPS or y + h > SH + EPS:
                issues.append(
                    f"slide {i}  BOUNDS     {str(name):28s} "
                    f"x={x:.2f} y={y:.2f} w={w:.2f} h={h:.2f} (slide {SW:.2f}x{SH:.2f})")
            if sh.has_chart:
                n_charts += 1
                blocks.append(("chart", r))
                continue
            if sh.has_table:
                n_tables += 1
                blocks.append(("table", r))
                continue
            if sh.has_text_frame and sh.text_frame.text.strip():
                need = est_height(sh.text_frame, w)
                if need > h * OVERFLOW_TOL:
                    issues.append(
                        f"slide {i}  OVERFLOW   needs {need:.2f}in in {h:.2f}in "
                        f"({need / h * 100:.0f}%)  '{sh.text_frame.text.strip()[:52]}'")
                texts.append((r, sh.text_frame.text.strip()[:40]))

        for tr, label in texts:
            for kind, br in blocks:
                if overlaps(tr, br):
                    issues.append(f"slide {i}  COLLISION  text over {kind}  '{label}'")

    print(f"deck   : {DECK}")
    print(f"slides : {len(prs.slides.__iter__.__self__._sldIdLst)}  "
          f"size {SW:.2f} x {SH:.2f} in")
    print(f"charts : {n_charts}   tables : {n_tables}")
    print(f"issues : {len(issues)}")
    for t in issues:
        print("  " + t)
    return 1 if issues else 0


if __name__ == "__main__":
    sys.exit(main())
