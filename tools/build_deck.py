#!/usr/bin/env python3
"""
Builds the 6-page Wayzyy moderation proposal deck.

Uses the Orange-and-Black template as the base so the theme, masters and fonts
carry over, then replaces all slides with six dense, purpose-built ones.

Design system extracted from the template:
    background  #F2EFEA   cream
    text        #000000   black
    accent      #CE4C32   burnt orange
    canvas      20.00 x 11.25 in  (16:9)
"""

import copy
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

SRC = "/Users/apple/Desktop/Orange and Black Modern Pitch Deck Presentation .pptx"
OUT = "/Users/apple/Desktop/Wayzyy_Moderation_Proposal_Deck.pptx"

CREAM = RGBColor(0xF2, 0xEF, 0xEA)
BLACK = RGBColor(0x00, 0x00, 0x00)
ORANGE = RGBColor(0xCE, 0x4C, 0x32)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GREY = RGBColor(0x5A, 0x55, 0x50)
LIGHT = RGBColor(0xE6, 0xE1, 0xD9)

FONT = "Calibri"
W, H = 20.0, 11.25
ML = 1.12
CW = W - 2 * ML  # content width 17.76

prs = Presentation(SRC)
prs.slide_width = Inches(W)
prs.slide_height = Inches(H)

# ---- wipe the template's slides, keep everything else -----------------------
sld_id_lst = prs.slides._sldIdLst
for sid in list(sld_id_lst):
    prs.part.drop_rel(sid.rId)
    sld_id_lst.remove(sid)

# pick the emptiest layout available
blank = min(prs.slide_layouts, key=lambda l: len(l.placeholders._element))


# ---- primitives -------------------------------------------------------------
def new_slide():
    s = prs.slides.add_slide(blank)
    for ph in list(s.placeholders):
        ph._element.getparent().remove(ph._element)
    bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(W), Inches(H))
    bg.fill.solid()
    bg.fill.fore_color.rgb = CREAM
    bg.line.fill.background()
    bg.shadow.inherit = False
    return s


def box(s, l, t, w, h, fill=None, line=None, lw=1.0, shape=MSO_SHAPE.RECTANGLE):
    sh = s.shapes.add_shape(shape, Inches(l), Inches(t), Inches(w), Inches(h))
    if fill is None:
        sh.fill.background()
    else:
        sh.fill.solid()
        sh.fill.fore_color.rgb = fill
    if line is None:
        sh.line.fill.background()
    else:
        sh.line.color.rgb = line
        sh.line.width = Pt(lw)
    sh.shadow.inherit = False
    return sh


def _est_h(items, w, size, space, line_spacing):
    """Estimate rendered height in inches. Matches the validator's model."""
    import math
    total = 0.0
    for item in items:
        over = {}
        if isinstance(item, tuple):
            item, over = item
        if not item:
            continue
        sz = over.get("size", size)
        ls = over.get("line_spacing", line_spacing)
        sa = over.get("space", space)
        char_w = 0.50 * sz / 72.0
        cpl = max(int(w / char_w), 1)
        lines = sum(max(1, math.ceil(len(seg) / cpl)) for seg in item.split("\n"))
        total += lines * (sz * ls * 1.20) / 72.0 + sa / 72.0
    return total


def text(s, l, t, w, h, runs, size=13, color=BLACK, bold=False, align=PP_ALIGN.LEFT,
         anchor=MSO_ANCHOR.TOP, space=4, line_spacing=0.95, italic=False,
         fit=True, min_size=10.5):
    """runs: str, or list of str, or list of (str, dict-overrides).

    With fit=True the type size is stepped down until the estimated rendered height
    fits the box. PowerPoint text boxes do not clip, so overflow silently collides
    with whatever sits below — this keeps the layout honest.
    """
    if isinstance(runs, str):
        runs = [runs]
    scale = 1.0
    if fit and h > 0:
        while scale > 0.55:
            if _est_h(runs, w, size * scale, space, line_spacing) <= h:
                break
            if size * scale <= min_size:
                break
            scale -= 0.04

    tb = s.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    first = True
    for item in runs:
        over = {}
        if isinstance(item, tuple):
            item, over = item
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.alignment = over.get("align", align)
        p.space_after = Pt(over.get("space", space))
        p.line_spacing = over.get("line_spacing", line_spacing)
        r = p.add_run()
        r.text = item
        f = r.font
        f.name = over.get("font", FONT)
        f.size = Pt(max(over.get("size", size) * scale, min_size))
        f.bold = over.get("bold", bold)
        f.italic = over.get("italic", italic)
        f.color.rgb = over.get("color", color)
    return tb


def header(s, num, title, kicker=None):
    text(s, ML, 0.78, 2.4, 1.0, num, size=46, bold=True, color=ORANGE)
    text(s, ML + 2.5, 0.86, 13.0, 0.9, title, size=40, bold=True, color=BLACK)
    if kicker:
        text(s, ML + 2.5, 1.86, 15.2, 0.4, kicker, size=15, color=GREY, italic=True)
    box(s, ML, 2.42, CW, 0.035, fill=ORANGE)


def footer(s, page):
    text(s, ML, 10.68, 9.0, 0.3, "Wayzyy Conversation Moderation — Technical Proposal",
         size=11, color=GREY)
    text(s, W - ML - 2.0, 10.68, 2.0, 0.3, f"{page} / 9", size=11, color=GREY,
         align=PP_ALIGN.RIGHT)


def card(s, l, t, w, h, title, lines, accent=ORANGE, tsize=15, bsize=12, fill=WHITE):
    box(s, l, t, w, h, fill=fill, line=LIGHT, lw=0.75)
    box(s, l, t, 0.055, h, fill=accent)
    tw = w - 0.5
    th = _est_h([title], tw, tsize, 0, 1.0) + 0.06
    text(s, l + 0.28, t + 0.17, tw, th, title, size=tsize, bold=True, color=BLACK,
         min_size=11.0)
    if lines:
        body_t = t + 0.17 + th + 0.16
        text(s, l + 0.28, body_t, tw, (t + h - 0.18) - body_t, lines, size=bsize,
             color=GREY, space=3, line_spacing=1.0, min_size=10.5)


def table(s, l, t, w, rows, col_w, header_size=12, body_size=11.5, row_h=0.34,
          head_h=0.38, zebra=True):
    """rows[0] is the header. col_w are relative weights."""
    n_r, n_c = len(rows), len(rows[0])
    total = sum(col_w)
    widths = [w * c / total for c in col_w]
    gt = s.shapes.add_table(n_r, n_c, Inches(l), Inches(t), Inches(w),
                            Inches(head_h + row_h * (n_r - 1)))
    tbl = gt.table
    tbl.first_row = True
    for i, cw in enumerate(widths):
        tbl.columns[i].width = Inches(cw)
    tbl.rows[0].height = Inches(head_h)
    for r in range(1, n_r):
        tbl.rows[r].height = Inches(row_h)
    for r, row in enumerate(rows):
        for c, val in enumerate(row):
            cell = tbl.cell(r, c)
            cell.margin_left = Inches(0.1)
            cell.margin_right = Inches(0.06)
            cell.margin_top = Inches(0.03)
            cell.margin_bottom = Inches(0.03)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            cell.fill.solid()
            if r == 0:
                cell.fill.fore_color.rgb = BLACK
            elif zebra and r % 2 == 0:
                cell.fill.fore_color.rgb = LIGHT
            else:
                cell.fill.fore_color.rgb = WHITE
            bold = False
            col = BLACK
            txt = val
            if isinstance(val, tuple):
                txt, opts = val
                bold = opts.get("bold", False)
                col = opts.get("color", BLACK)
            tf = cell.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.LEFT
            r_ = p.add_run()
            r_.text = str(txt)
            f = r_.font
            f.name = FONT
            f.size = Pt(header_size if r == 0 else body_size)
            f.bold = True if r == 0 else bold
            f.color.rgb = WHITE if r == 0 else col
    return gt


MONO = "Consolas"

# ---- charts ----------------------------------------------------------------
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION, XL_LABEL_POSITION


def _chart_font(chart, size=10.0):
    chart.font.size = Pt(size)
    chart.font.name = FONT
    chart.font.color.rgb = GREY


def column_chart(s, l, t, w, h, cats, series, colors, number_format='0.0"%"',
                 legend=True, label_size=9.5, max_scale=None, gap=90, overlap=-15):
    """Clustered column chart on a white card, styled to the deck palette."""
    box(s, l, t, w, h, fill=WHITE, line=LIGHT, lw=0.75)
    cd = CategoryChartData()
    cd.categories = cats
    for name, vals in series:
        cd.add_series(name, vals)
    gf = s.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED,
                            Inches(l + 0.10), Inches(t + 0.10),
                            Inches(w - 0.20), Inches(h - 0.20), cd)
    ch = gf.chart
    _chart_font(ch, 10.0)

    ch.has_legend = legend and len(series) > 1
    if ch.has_legend:
        ch.legend.position = XL_LEGEND_POSITION.TOP
        ch.legend.include_in_layout = False
        ch.legend.font.size = Pt(9.5)

    plot = ch.plots[0]
    plot.gap_width = gap
    plot.overlap = overlap
    plot.has_data_labels = True
    dl = plot.data_labels
    dl.number_format = number_format
    dl.number_format_is_linked = False
    dl.font.size = Pt(label_size)
    dl.font.bold = True
    dl.font.color.rgb = BLACK
    try:
        dl.position = XL_LABEL_POSITION.OUTSIDE_END
    except Exception:
        pass

    for i, ser in enumerate(ch.series):
        ser.format.fill.solid()
        ser.format.fill.fore_color.rgb = colors[i % len(colors)]
        ser.format.line.fill.background()

    va = ch.value_axis
    va.has_major_gridlines = False
    va.visible = False
    if max_scale is not None:
        va.maximum_scale = max_scale
    va.minimum_scale = 0

    ca = ch.category_axis
    ca.has_major_gridlines = False
    ca.tick_labels.font.size = Pt(9.5)
    ca.tick_labels.font.color.rgb = BLACK
    ca.format.line.color.rgb = LIGHT
    return ch


def doughnut_chart(s, l, t, w, h, cats, vals, colors, number_format='0.0"%"',
                   label_size=10.0):
    box(s, l, t, w, h, fill=WHITE, line=LIGHT, lw=0.75)
    cd = CategoryChartData()
    cd.categories = cats
    cd.add_series("s", vals)
    gf = s.shapes.add_chart(XL_CHART_TYPE.DOUGHNUT,
                            Inches(l + 0.10), Inches(t + 0.10),
                            Inches(w - 0.20), Inches(h - 0.20), cd)
    ch = gf.chart
    _chart_font(ch, 10.0)
    ch.has_legend = True
    ch.legend.position = XL_LEGEND_POSITION.BOTTOM
    ch.legend.include_in_layout = False
    ch.legend.font.size = Pt(9.5)
    ch.legend.font.color.rgb = BLACK

    plot = ch.plots[0]
    plot.has_data_labels = True
    dl = plot.data_labels
    dl.number_format = number_format
    dl.number_format_is_linked = False
    dl.font.size = Pt(label_size)
    dl.font.bold = True
    dl.font.color.rgb = WHITE

    pts = ch.series[0].points
    for i, pt in enumerate(pts):
        pt.format.fill.solid()
        pt.format.fill.fore_color.rgb = colors[i % len(colors)]
        pt.format.line.color.rgb = WHITE
        pt.format.line.width = Pt(1.5)
    return ch


def chart_label(s, l, t, w, label):
    text(s, l, t, w, 0.24, label, size=9.5, bold=True, color=ORANGE, fit=False)


def field(s, l, t, w, h, label, body, bsize=10.8, mono=False, color=GREY):
    """A labelled field inside a layer panel."""
    text(s, l, t, w, 0.24, label, size=9.2, bold=True, color=ORANGE, fit=False)
    text(s, l, t + 0.26, w, h - 0.26, body, size=bsize, color=color,
         space=2, line_spacing=1.02)


def layer_panel(s, t, h, tag, name, is_what, does, made_of, example,
                accent=ORANGE, emphasis=False):
    """Full-width panel describing one layer: what it is / does / made of / example.

    Field heights are allocated from the estimated height of the copy rather than by a
    fixed fraction, so a short "what it is" hands its slack to "what it does" instead of
    forcing the auto-fit down to its floor.
    """
    l, w = ML, CW
    box(s, l, t, w, h, fill=WHITE, line=(ORANGE if emphasis else LIGHT),
        lw=1.25 if emphasis else 0.75)
    box(s, l, t, 0.055, h, fill=accent)

    text(s, l + 0.24, t + 0.19, 2.10, 0.26, tag, size=10.5, bold=True, color=accent)
    text(s, l + 0.24, t + 0.47, 2.10, 0.82, name, size=16, bold=True, color=BLACK,
         line_spacing=0.92)

    # column A carries the most prose, so it gets the most width
    xa, wa = l + 2.52, 6.15
    xb, wb = l + 8.92, 4.15
    xc, wc = l + 13.42, 4.05
    top = t + 0.19
    inner_h = h - 0.36

    ha = min(_est_h([is_what], wa, 10.8, 0, 1.02) + 0.30, inner_h * 0.52)
    field(s, xa, top, wa, ha, "WHAT IT IS", is_what, bsize=10.8, color=BLACK)
    field(s, xa, top + ha + 0.08, wa, inner_h - ha - 0.08, "WHAT IT DOES", does)

    field(s, xb, top, wb, inner_h, "MADE OF", made_of)

    box(s, xc - 0.14, t + 0.15, wc + 0.28, inner_h + 0.08, fill=CREAM)
    text(s, xc, t + 0.22, wc, 0.24, "EXAMPLE", size=9.2, bold=True, color=ORANGE, fit=False)
    text(s, xc, t + 0.48, wc, inner_h - 0.31, example, size=10.4, color=BLACK,
         space=2, line_spacing=1.05, min_size=9.6)


def chain(s, t, steps, active):
    """Horizontal cascade chain; `active` is a set of indices to highlight."""
    n = len(steps)
    gap, aw = 0.30, 0.20
    bw = (CW - (n - 1) * (gap + aw)) / n
    x = ML
    for i, (lab, sub) in enumerate(steps):
        on = i in active
        box(s, x, t, bw, 0.62, fill=(BLACK if on else WHITE),
            line=(None if on else LIGHT), lw=0.75)
        text(s, x, t + 0.10, bw, 0.24, lab, size=11.5, bold=True,
             color=(WHITE if on else GREY), align=PP_ALIGN.CENTER, fit=False)
        text(s, x, t + 0.34, bw, 0.22, sub, size=9.5,
             color=(CREAM if on else GREY), align=PP_ALIGN.CENTER, fit=False)
        x += bw
        if i < n - 1:
            box(s, x + gap / 2 - 0.02, t + 0.22, aw, 0.18,
                fill=ORANGE, shape=MSO_SHAPE.RIGHT_ARROW)
            x += gap + aw


def metric(s, l, t, w, value, label, vsize=30):
    text(s, l, t, w, 0.5, value, size=vsize, bold=True, color=ORANGE)
    text(s, l, t + 0.52, w, 0.6, label, size=11.5, color=GREY, line_spacing=1.0)


def chip(s, l, t, w, h, label, fill=BLACK, fg=WHITE, size=11.5, bold=True):
    box(s, l, t, w, h, fill=fill)
    text(s, l, t + (h - 0.22) / 2 - 0.03, w, 0.3, label, size=size, bold=bold,
         color=fg, align=PP_ALIGN.CENTER)


def arrow_down(s, cx, t, h):
    a = box(s, cx - 0.09, t, 0.18, h, fill=ORANGE, shape=MSO_SHAPE.DOWN_ARROW)
    return a


# ============================================================================
# SLIDE 1 — Title and the answer up front
# ============================================================================
s = new_slide()
box(s, 0, 0, 0.42, H, fill=BLACK)
box(s, W - 5.7, 0, 5.7, H, fill=ORANGE)

text(s, ML + 0.5, 1.15, 5.0, 0.35, "TECHNICAL PROPOSAL", size=13, bold=True, color=ORANGE)
text(s, ML + 0.5, 1.95, 12.0, 2.9,
     ["Conversation", "Moderation for Wayzyy"],
     size=62, bold=True, color=BLACK, space=0, line_spacing=0.92)
box(s, ML + 0.5, 4.62, 4.2, 0.05, fill=ORANGE)

text(s, ML + 0.5, 5.05, 12.3, 1.5,
     ["Off-platform contact sharing and conversation safety — threats, harassment,",
      "blackmail, sexual harassment and phishing — solved as ONE service that reads",
      "each message once and returns one verdict."],
     size=17, color=GREY, space=3, line_spacing=1.05)

# headline metric strip
metrics = [
    ("96.0%", "contact recall\nat 100% precision"),
    ("0.00%", "false positives\non innocent traffic"),
    ("3.45 ms", "p50 write-path\nlatency"),
    ("10/10", "Wayzyy official\ntest cases"),
    ("18/18", "unseen multilingual\nsafety cases"),
]
box(s, ML + 0.5, 6.75, 12.3, 0.035, fill=LIGHT)
for i, (v, lab) in enumerate(metrics):
    metric(s, ML + 0.5 + i * 2.48, 7.05, 2.4, v, lab, vsize=27)

box(s, ML + 0.5, 8.72, 12.3, 0.035, fill=LIGHT)
text(s, ML + 0.5, 9.0, 12.3, 0.9,
     [("Deterministic tiers decide 94.7% of messages with no model and no network. "
       "A language model is consulted only for the ~5% that are genuinely ambiguous.",
       {"size": 14, "color": BLACK, "bold": True})],
     line_spacing=1.1)

# right panel
text(s, W - 5.7 + 0.75, 1.6, 4.2, 0.4, "AT A GLANCE", size=13, bold=True, color=WHITE)
box(s, W - 5.7 + 0.75, 2.12, 1.8, 0.04, fill=WHITE)
glance = [
    ("Architecture", "4-tier cascade, one microservice"),
    ("Write path", "in-process, no network hop"),
    ("Model tier", "off the write path, async"),
    ("Throughput", "285 msg/s per core"),
    ("Concurrency", "~45,000 active senders / instance"),
    ("Cost", "$0.035 per million (tiers 1–2)"),
    ("", "$0.74 per 100k (model tier)"),
    ("Contact failure", "fails OPEN — chat never stalls"),
    ("Safety failure", "fails CLOSED — held for review"),
    ("Portability", "Swift, Linux-ready, vendor-agnostic"),
]
y = 2.55
for k, v in glance:
    if k:
        text(s, W - 5.7 + 0.75, y, 4.3, 0.26, k, size=11.5, bold=True, color=BLACK)
        y += 0.27
    text(s, W - 5.7 + 0.75, y, 4.3, 0.3, v, size=12.5, color=WHITE, line_spacing=1.0)
    y += 0.42

footer(s, 1)

# ============================================================================
# SLIDE 2 — Architecture 1 of 3: the deterministic path
# ============================================================================
s = new_slide()
header(s, "01.", "Architecture 1 of 3 — the deterministic path",
       "Tiers 0 and 1 run on 100% of messages, cost nothing but CPU, and decide 94.7% of all traffic on their own.")

chain(s, 2.60, [("TIER 0", "canonicalise"), ("TIER 1", "extract + validate"),
                ("TIER 2", "retrieve"), ("LAYER 3", "classify"),
                ("FUSE", "score + policy"), ("ROUTER", "ask or not"),
                ("TIER 3", "reason")], active={0, 1})

text(s, ML, 3.28, CW, 0.48,
     [("Why contact and safety are one service:  ", {"bold": True, "size": 11.5}),
      ("a message can carry a phone number AND a threat. With two systems the first to fire wins and the "
       "other finding is lost — reproduced in testing, where an extortion message scored as a contact "
       "issue and its threat went unexamined.", {"size": 11, "color": GREY})], line_spacing=1.0)

layer_panel(
    s, 3.86, 3.30,
    "TIER 0  ·  0.5 ms  ·  100%", "CANONICALISE",
    "A normalisation stage that rewrites every message into 13 clean forms before anything inspects it.",
    ["Undoes disguise generically rather than chasing each trick, so it generalises to evasions nobody "
     "has seen yet. Keeps a character map back to the original, which is why redaction is exact and "
     "every finding is explainable.",
     "One view exists purely to kill false positives: digitsMasked strips prices, counts, times and "
     "reference numbers before phone matching runs.",
     "Must run BEFORE any learned component. Guardrail models are measurably fragile to character "
     "mutation, and this layer recovers most of that loss for free."],
    ["13 views — raw · base · alpha · compact · alphaCompact · digits · digitsMasked · digitsReversed · "
     "separators · separatorsAlt · acrostic · compactDigits · romanDigits",
     "NFKC compatibility fold · homoglyph table · zero-width and TAG-block decode · leet map · repeat "
     "collapse · number-word expansion (EN + Hindi) · separator words (at / dot / underscore) · Roman "
     "numerals",
     "Two platform-specific Unicode functions, isolated — so the same engine runs on Linux."],
    [("in    Са1l Ме оn nine eight 7 six zero", {"bold": True, "size": 10.4, "color": BLACK}),
     ("      Cyrillic С а М о — looks Latin", {"size": 9.6, "color": GREY}),
     ("base  call me on nine eight 7 six zero", {"size": 10.4}),
     ("dig   98760", {"size": 10.4}),
     ("", {"size": 5}),
     ("in    Total is 12,500 for 3 nights", {"bold": True, "size": 10.4, "color": BLACK}),
     ("dig          125003   ← phone-shaped", {"size": 10.4}),
     ("digMasked    (empty)  →  ALLOW", {"size": 10.4, "color": ORANGE, "bold": True})])

layer_panel(
    s, 7.28, 3.30,
    "TIER 1  ·  1.5 ms  ·  100%", "EXTRACT + VALIDATE",
    "Deterministic detection that proposes candidates, then validates them against real-world rules "
    "before believing them.",
    ["Runs two branches that never gate each other — contact extraction and the safety floor — so a "
     "phone number cannot satisfy the engine into skipping a threat.",
     "Validation is where precision comes from: extraction proposes, published numbering plans decide. "
     "That is what gives 100% precision rather than a tuned threshold.",
     "A rolling window per sender re-analyses the concatenation, so a number split across five messages "
     "is still recovered. The window is consumed once it fires."],
    ["Extractors — phone · e-mail · URL · social handle · payment VPA · platform steering · bare "
     "identifier",
     "Numbering-plan validators (Indian mobile, NANP, E.164) · self-label rule for plan-invalid numbers",
     "Safety floor — 259 terms across threat / harassment / coercion / scam / sexual / self-harm / "
     "profanity, English + romanised Hindi · target-aware profanity resolver · URL parser with "
     "registrable-domain compare",
     "Decoders — Base64 · hex · ROT13 ×25 · binary · Morse · NATO · percent-encoding"],
    [("987 654 3210", {"bold": True, "size": 10.4, "color": BLACK}),
     ("  → valid Indian mobile  →  MASK", {"size": 10.4}),
     ("12500", {"bold": True, "size": 10.4, "color": BLACK}),
     ("  → no plan matches  →  not a phone", {"size": 10.4}),
     ("1234567891", {"bold": True, "size": 10.4, "color": BLACK}),
     ("  → rejected, no plan allows a lead 1", {"size": 10.4}),
     ("", {"size": 4}),
     ("mera number ek do teen chaar…", {"bold": True, "size": 10.4, "color": BLACK}),
     ("  → plan-invalid, but the sender said", {"size": 10.4}),
     ("    it IS their number → share, 0.86", {"size": 10.4, "color": ORANGE, "bold": True})])

footer(s, 2)

# ============================================================================
# SLIDE 3 — Architecture 2 of 3: the learned layer
# ============================================================================
s = new_slide()
header(s, "02.", "Architecture 2 of 3 — the learned layer",
       "Where meaning is judged rather than matched. Retrieval handles what has no payload to extract; the classifier makes coverage language-general.")

chain(s, 2.60, [("TIER 0", "canonicalise"), ("TIER 1", "extract + validate"),
                ("TIER 2", "retrieve"), ("LAYER 3", "classify"),
                ("FUSE", "score + policy"), ("ROUTER", "ask or not"),
                ("TIER 3", "reason")], active={2, 3})

text(s, ML, 3.28, CW, 0.44,
     [("These two layers are why coverage is not limited to the words we predicted.  ",
       {"bold": True, "size": 11.5}),
      ("Tier 2 generalises by example and explains itself; Layer 3 generalises across languages and "
       "gives a calibrated probability per category. Both run before any money is spent.",
       {"size": 11, "color": GREY})], line_spacing=1.0)

layer_panel(
    s, 3.82, 3.32,
    "TIER 2  ·  1.2 ms  ·  ~all", "RETRIEVE",
    "Nearest-neighbour retrieval over labelled example messages — a non-parametric classifier, used "
    "instead of a trained one.",
    ["Catches meaning with no payload to extract — referential sharing (\u201cmy number is in my profile "
     "bio\u201d has no digits at all), and abuse phrased in words no list contains.",
     "Acts only when the match beats the nearest INNOCENT example by a margin, because ordinary travel "
     "chat is mildly similar to everything.",
     "Chosen over a trained classifier for three reasons: it generalises by adding examples rather than "
     "retraining, it quotes the neighbour it matched so an appeal reviewer sees evidence, and it is "
     "deterministic and cacheable."],
    ["336 anchors — 98 contact · 113 safety · 125 innocent",
     "TWO interchangeable backends behind one protocol.  Lexical: hashed n-gram TF-IDF, word "
     "uni/bigrams + char 4-grams, 2^18 space, bars 0.28 / 0.06.  Embedding: any OpenAI-compatible "
     "/v1/embeddings endpoint, bars 0.60 / 0.05.",
     "Each space carries its OWN calibrated bars and a vector travels with its space — cosine is not "
     "comparable across spaces. Lexical is the default on the measurement: equal recall, 20% faster.",
     "Contact and safety pools never merged, so one family cannot mask the other"],
    [("my number is in my profile bio", {"bold": True, "size": 10.4, "color": BLACK}),
     ("no digits at all — Tier 1 finds nothing", {"size": 9.6, "color": GREY}),
     ("nearest contact anchor   0.71", {"size": 10.4}),
     ("nearest innocent anchor  0.14", {"size": 10.4}),
     ("margin 0.57  →  MASK", {"size": 10.4, "color": ORANGE, "bold": True}),
     ("verdict quotes the anchor it matched", {"size": 9.6, "color": GREY})])

layer_panel(
    s, 7.26, 3.32,
    "LAYER 3  ·  ~8 ms  ·  100%", "CLASSIFY",
    "A quantised multilingual multi-label classifier, INT8 on CPU, run on every message.",
    ["Turns safety coverage from lexicon-general into language-general. One model spans dozens of "
     "languages including code-mixed Hinglish, which no wordlist can.",
     "Emits a calibrated probability per category, so enforcement thresholds become an explicit, "
     "revisitable trade instead of a guess. Three bands: enforce · route · allow.",
     "The legitimate-complaint head is the unusual one and it is deliberate — it gives the system a "
     "positive way to say \u201cthis is an angry customer exercising a right\u201d rather than inferring it "
     "from the absence of other signals."],
    ["XLM-R class encoder, INT8 via ONNX Runtime, CPU only — no GPU, no network",
     "Seven INDEPENDENT sigmoid heads — threat · harassment · sexual · self-harm · coercion · scam · "
     "legitimate-complaint.  Multi-label rather than softmax, because a message can be two things at "
     "once and softmax would quietly lose one.",
     "~$0.08 per million messages · bands calibrated by sweeping the innocent suite with false-positive "
     "rate pinned at target, never by intuition"],
    [("tum bilkul ghatiya insaan ho", {"bold": True, "size": 10.4, "color": BLACK}),
     ("romanised Hindi, no listed swear word", {"size": 9.6, "color": GREY}),
     ("harassment            0.91", {"size": 10.4}),
     ("threat                0.06", {"size": 10.4}),
     ("legitimate-complaint  0.04", {"size": 10.4}),
     ("→ enforce band  →  WARN", {"size": 10.4, "color": ORANGE, "bold": True})],
    emphasis=True)

footer(s, 3)

# ============================================================================
# SLIDE 4 — Architecture 3 of 3: decide, route, reason
# ============================================================================
s = new_slide()
header(s, "03.", "Architecture 3 of 3 — decide, route, reason",
       "Where the action is chosen, how we decide what is worth paying for, and why the only paid tier sits after delivery.")

chain(s, 2.60, [("TIER 0", "canonicalise"), ("TIER 1", "extract + validate"),
                ("TIER 2", "retrieve"), ("LAYER 3", "classify"),
                ("FUSE", "score + policy"), ("ROUTER", "ask or not"),
                ("TIER 3", "reason")], active={4, 5, 6})

layer_panel(
    s, 3.34, 2.20,
    "FUSION + POLICY  ·  0.1 ms", "SCORE + DECIDE",
    "The scoring and decision layer — one score from every finding, then one action.",
    ["Moves thresholds by who is speaking and where the booking is, rather than one global bar that "
     "over-blocks trusted hosts and under-blocks new accounts.",
     "Prices obfuscation effort as intent — nobody writes a number in Cyrillic by accident.",
     "Safety categories bypass thresholds entirely."],
    ["Weighted feature fusion · effort scored per detection from the view that produced it",
     "Risk-adaptive thresholds — hint 0.22 / mask 0.38 / withhold 0.60, offset by trust tier, booking "
     "stage and prior violations · explicit-identifier floor, so a literally stated handle masks "
     "whatever the thresholds say",
     "Six-rung action ladder — allow · hint · mask · warn · block · review"],
    [("nine eight 7 six zero 1 2 3 4", {"bold": True, "size": 10.4, "color": BLACK}),
     ("score       0.91", {"size": 10.4}),
     ("effort      3  (number-words)", {"size": 10.4}),
     ("sender      verified · inquiry stage", {"size": 10.4}),
     ("threshold   withhold 0.60", {"size": 10.4}),
     ("→ MASK, redact characters 34–62", {"size": 10.4, "color": ORANGE, "bold": True})])

layer_panel(
    s, 5.66, 2.20,
    "ROUTER  ·  0.2 ms", "ASK OR NOT",
    "The gate that decides what is worth paying for. It can never enforce — its only power is to spend "
    "one model call.",
    ["Structure, not vocabulary — is a person addressed, does this resemble anything people normally "
     "write, is there a conditional demand, is the score in the grey band.",
     "The inversion that lets a finite system cover an infinite problem: a word missing from a lexicon "
     "now means ASK, not ALLOW."],
    ["10 structural signals",
     "~29 second-person pronouns across English, Hinglish, Devanagari and Cyrillic — matched on the "
     "ORIGINAL text, because the homoglyph fold rewrites Cyrillic words",
     "Novelty distance from the innocent anchor corpus · conditional connectives (or else · unless you · "
     "warna · иначе) · uncertainty band 0.10–0.62"],
    [("tujhe maar dunga", {"bold": True, "size": 10.4, "color": BLACK}),
     ("romanised Hindi — \u201cI will kill you\u201d", {"size": 9.6, "color": GREY}),
     ("addresses a person?  tujhe   → YES", {"size": 10.4}),
     ("resembles normal?    0.00    → NO", {"size": 10.4}),
     ("→ route to Tier 3", {"size": 10.4, "color": ORANGE, "bold": True}),
     ("Nothing in the code knows what those", {"size": 10.4, "color": ORANGE, "bold": True}),
     ("words mean.", {"size": 10.4, "color": ORANGE, "bold": True})])

chip(s, ML, 7.98, CW, 0.42,
     "VERDICT RETURNED  ·  p50 3.45 ms  ·  MESSAGE DELIVERED  —  everything below this line is asynchronous",
     fill=BLACK, size=12.5)

layer_panel(
    s, 8.52, 2.06,
    "TIER 3  ·  1.7 s  ·  5.3%", "REASON",
    "A small language model as judge, deliberately off the write path.",
    ["Answers the one genuinely open-ended question: natural language or constructed evasion, and is "
     "this implication a threat.",
     "Then revises — retro-mask, withhold, clear or hold for review. Contact fails open; safety fails "
     "closed."],
    ["gpt-oss-20b at LOW reasoning effort — measured faster, cheaper AND more accurate than the same "
     "model at default effort",
     "OpenAI chat-completions shape, so hosted APIs, vLLM and Ollama are a base URL and a model name · "
     "8-message window · injection fencing · strict 8-decision JSON schema · pooled lanes with "
     "per-model quota isolation · rate buckets and a self-healing breaker"],
    [("add one to every digit of the", {"bold": True, "size": 10.4, "color": BLACK}),
     ("booking ref — that's my mobile", {"bold": True, "size": 10.4, "color": BLACK}),
     ("the number is not in the text —", {"size": 9.6, "color": GREY}),
     ("it is a RULE to construct one", {"size": 9.6, "color": GREY}),
     ("delivered 3 ms · judged 1.7 s later", {"size": 10.4}),
     ("0.96 → retro-masked", {"size": 10.4, "color": ORANGE, "bold": True})],
    emphasis=True)

footer(s, 4)

# ============================================================================
# SLIDE 5 — Model selection evidence
# ============================================================================
s = new_slide()
header(s, "04.", "Model selection — the evidence",
       "Every learned component was chosen by measurement, and two candidates were rejected on the numbers. This is that record.")

LW = 11.40
chart_label(s, ML, 2.70, LW, "TIER 2 REPRESENTATION  —  LEXICAL vs EMBEDDINGS, SAME 480-CASE CORPUS")
table(s, ML, 3.00, LW, [
    ["", "A · LEXICAL  (default)", "B · EMBEDDINGS  (optional)"],
    ["Representation", "Hashed n-gram TF-IDF", "Dense sentence vectors"],
    ["Features", "Word uni + bigrams, char 4-grams, 2^18", "Model-defined"],
    ["Default model", "— (no model, no network)", "nomic-embed-text"],
    ["Similarity / margin bar", "0.28 / 0.06", "0.60 / 0.05"],
    ["Confidence ceiling", "0.62", "0.90"],
    [("Measured recall @ 0.00% FPR", {"bold": True}), ("85.9%", {"bold": True}), ("85.9%", {"bold": True})],
    [("Measured latency", {"bold": True}), ("2.13 ms", {"bold": True}),
     ("2.67 ms", {"bold": True, "color": ORANGE})],
], col_w=[3.2, 4.1, 4.1], row_h=0.34, head_h=0.34, body_size=11)

text(s, ML, 5.86, LW, 0.46,
     [("Equal recall, 25% more latency, plus an operational dependency.", {"bold": True, "size": 11.5, "color": ORANGE, "space": 1}),
      ("So lexical is the default and embeddings are a configuration option. We are not paying for a "
       "model that does not win.", {"size": 11, "color": GREY, "space": 1})], line_spacing=1.02)

chart_label(s, ML, 6.46, LW, "TIER 3 SLM BENCH  —  WHY A SMALL MODEL AT LOW REASONING EFFORT")
table(s, ML, 6.76, LW, [
    ["MODEL", "OUTCOME", "LATENCY", "COST / CALL", "WHAT THE MEASUREMENT SHOWED"],
    [("gpt-oss-20b · low effort", {"bold": True}), ("SELECTED", {"bold": True}), "586 ms",
     ("$0.000075", {"bold": True}), ("109 output tokens · 4 of 4 correct", {"bold": True})],
    ["gpt-oss-20b · default effort", "rejected", "985 ms", "$0.000075",
     "414 tokens · 3 of 4 — extra thinking bought a false positive"],
    ["llama-3.1-8b-instant", "rejected", "fast", "$0.000029",
     "Cheapest of all, but fails the word-length coincidence case"],
    ["qwen3-32b", "rejected", "3.1 s", "$0.000996", "13x the cost, no accuracy gain"],
    ["llama3.1:8b · local Ollama", "self-host lane", "~1.05 s", "electricity",
     "96.0% full-cascade recall · ~0.95 calls/s on an M1 Pro"],
], col_w=[3.0, 1.5, 1.2, 1.4, 4.3], row_h=0.40, head_h=0.34, body_size=11)

text(s, ML, 9.36, LW, 0.62,
     [("The counter-intuitive row is the second one.", {"bold": True, "size": 11.5, "color": ORANGE, "space": 1}),
      ("The same model, thinking harder, was slower, no cheaper, and less accurate. It broke on "
       "\u201cBreakfast is 250 per person, dinner around 600 for 2\u201d — whose word lengths validate as a real "
       "Indian mobile. Extra deliberation talked it into treating a coincidence as meaningful. More "
       "thinking is not free precision.", {"size": 11, "color": GREY, "space": 1})], line_spacing=1.02)

# right column
rx = ML + 11.95
card(s, rx, 2.70, 5.81, 3.34, "Cross-lingual embeddings — measured, then rejected",
     ["We tested whether multilingual embeddings could carry safety across languages. That would have "
      "been the cheap answer, so it was worth measuring before designing around it.",
      "nomic-embed-text — separated only 7 of 10 cases. Hinglish innocents overlapped threats outright.",
      "bge-m3 — 9 of 10 by margin SIGN, but threat margins ran as low as +0.002 against innocent "
      "\u22120.000. A margin of two thousandths is not a decision boundary, it is noise.",
      "Conclusion: usable for routing, not for enforcement. Multilingual coverage is therefore carried "
      "by the router's structural signals and by Layer 3 — not by embedding similarity.",
      "Declining to ship something because the measurement said no is what keeps the rest of these "
      "numbers meaningful."],
     bsize=11.5)

card(s, rx, 6.16, 5.81, 3.82, "The bug that shaped this layer",
     ["Cosine similarity is NOT comparable across vector spaces. Hashed n-grams leave unrelated text "
      "near zero, so 0.28 is a meaningful bar. Sentence embeddings put any two English sentences between "
      "0.45 and 1.0 — against that, 0.28 is below the floor and the gate is wide open.",
      "Our embedding backend degrades to lexical vectors when its endpoint is slow. For a period it "
      "returned them BARE — so lexical vectors were compared against embedded anchors and judged by "
      "embedding bars. Recall fell BELOW the plain lexical backend it was meant to beat, while every "
      "individual piece still looked correct in isolation.",
      "Fix: a vector always travels with the space that produced it, and each space carries its own bars "
      "and ceiling. The type system now makes the mistake unspellable.",
      "This is also why the bars are not global constants. Swept on the 480-case corpus with FPR pinned "
      "at zero, the embedding backend spans 85.4% to 88.9% recall on threshold choice alone — carrying "
      "the lexical numbers across would have cost 3.5 points silently."],
     bsize=11.5)

footer(s, 5)

# ============================================================================
# SLIDE 6 — Model and regex choices
# ============================================================================
s = new_slide()
header(s, "05.", "Model and regex choices",
       "Every question is routed to the cheapest mechanism that can answer it correctly. Nothing more powerful, nothing less.")

table(s, ML, 2.72, 17.76, [
    ["QUESTION", "MECHANISM", "WHY THIS AND NOT SOMETHING ELSE"],
    [("Is this a phone number?", {"bold": True}), "Deterministic extraction + validation against real numbering plans",
     "Closed-form. Numbering plans are published and finite. Validation is what delivers 100% precision — 987 654 3210 validates, 12500 does not, 1234567891 does not."],
    [("Is it disguised?", {"bold": True}), "13 canonical views, then one matcher",
     "One matcher per evasion is an infinite game. Undoing the disguise generalises to techniques nobody has seen. Ca1l Ме (Cyrillic) → call me."],
    [("Is this number legitimate?", {"bold": True}), "Numeric-context masking (view 7)",
     "The false-positive killer. 12,500 for 3 nights, ref WZ4471829 → masked to empty → ALLOW. A system that flags this makes hosts abandon chat."],
    [("Is the contact hidden by reference?", {"bold": True}), "Retrieval over 98 labelled anchors + margin",
     'No payload exists to extract. "my number is in my profile bio" has no digits. Retrieval generalises without code and quotes the anchor it matched.'],
    [("Is this host ours?", {"bold": True}), "URL parsing + registrable-domain compare",
     "Arithmetic, not judgement. Never a regex — regex read wayzyy.com@evil.example as our own domain. Paying a model for a deterministic answer is the mistake to avoid."],
    [("Is this person being abused?", {"bold": True}), "Target rule + 259-term fast path + retrieval",
     '"this place is shit" is a review; "you are shit" is abuse. The target decides, not the vocabulary. Wordlists are a fast path, never the mechanism.'],
    [("Is this abuse in a language we never listed?", {"bold": True}), "Vocabulary-free routing signals",
     "Direct address + novelty + conditional connective. Closed classes, stable for centuries. Routes to a model that already knows the language."],
    [("Is this a veiled threat or blackmail?", {"bold": True}), "LLM judge, small model, low reasoning effort",
     "Implication is not enumerable. No individual word in blackmail is hostile. This is the only genuinely open-ended question in the system."],
], col_w=[3.0, 4.3, 10.4], row_h=0.62, head_h=0.36, header_size=11.5, body_size=11)

# lower panels
ly = 8.30
card(s, ML, ly, 5.68, 2.28, "Where regex belongs — and where it does not",
     ["USE for closed-form structure: digit runs, e-mail and VPA shape, separator resolution, "
      "conditional connectives.",
      "DO NOT USE for threats. Hostility has no closed form. An official test case — \"give me "
      "full refund or i will trash the place\" — scored 0.028 and was allowed, because the "
      "extortion pattern required the verb to be report / review / rate / complain."],
     bsize=12)

card(s, ML + 6.04, ly, 5.68, 2.28, "Retrieval rather than a trained classifier",
     ["Generalises without code — a new phrasing lands near an existing anchor. Adding coverage "
      "is adding examples, not shipping a release.",
      "Explains itself — the verdict quotes the anchor it matched. An appeal reviewer can "
      "adjudicate that; a bare score is not reviewable.",
      "Deterministic and cacheable — no sampling variance, no prompt drift."],
     bsize=12)

card(s, ML + 12.08, ly, 5.68, 2.28, "Model choice — measured, not assumed",
     ["A 20B model at LOW reasoning effort: 109 output tokens, 586 ms, 4/4 correct. The SAME model "
      "at default effort: 414 tokens, 985 ms, 3/4. Cheaper, faster AND more accurate.",
      "Extra deliberation talked it into treating a coincidental decode as meaningful.",
      "The interface is the OpenAI chat-completions shape, so hosted APIs, vLLM, Ollama and "
      "llama.cpp are a base URL and a model name — no vendor lock-in."],
     bsize=12)

footer(s, 6)

# ============================================================================
# SLIDE 7 — Trade-offs: latency vs accuracy vs cost
# ============================================================================
s = new_slide()
header(s, "06.", "Trade-offs: latency vs accuracy vs cost",
       "We built and measured all three approaches on the same 480-case corpus. The cascade is not a compromise between them — it wins on every axis at once.")

LW = 11.40
chart_label(s, ML, 2.70, LW, "THE THREE APPROACHES  —  SAME 480-CASE CORPUS, SAME MACHINE")
table(s, ML, 3.02, LW, [
    ["APPROACH", "RECALL", "FALSE POS.", "LATENCY", "COST /100k"],
    ["Regex and keyword lists only", "34.3%", ("4.88%", {"color": ORANGE, "bold": True}), "0.01 ms", "$0"],
    ["A language model on every message", "88.1%", ("13.2%", {"color": ORANGE, "bold": True}), "~2,400 ms", "$7.50"],
    [("This proposal — 4-tier cascade", {"bold": True}), ("96.0%", {"bold": True}),
     ("0.00%", {"bold": True}), ("2.3 ms", {"bold": True}), ("$0.74", {"bold": True})],
], col_w=[4.6, 1.6, 1.7, 1.7, 1.8], row_h=0.44, head_h=0.36, body_size=12)

text(s, ML, 4.78, LW, 0.62,
     [("Read the middle row twice.", {"bold": True, "size": 12.5, "color": ORANGE, "space": 1}),
      ("A language model on every message is worse on accuracy, ~900x slower and ~10x the cost. It "
       "over-flags because a model asked \u201cis this suspicious?\u201d with no deterministic grounding "
       "always finds something to worry about.", {"size": 11.5, "color": GREY, "space": 1})],
     line_spacing=1.02)

chart_label(s, ML, 5.42, 5.50, "ACCURACY  —  RECALL AGAINST FALSE POSITIVES")
column_chart(s, ML, 5.70, 5.50, 2.20,
             ["Regex only", "LLM on every msg", "This proposal"],
             [("Recall %", (34.3, 88.1, 96.0)),
              ("False positives %", (4.88, 13.2, 0.0))],
             colors=[BLACK, ORANGE], number_format='0.0"%"', max_scale=100)

chart_label(s, ML + 5.90, 5.42, 5.50, "WHAT EACH TIER BUYS  —  RECALL vs SPEND")
column_chart(s, ML + 5.90, 5.70, 5.50, 2.20,
             ["Tier 1", "T1+2 lexical", "T1+2 embed", "T1+2+3"],
             [("Recall %", (81.9, 85.9, 85.9, 96.0)),
              ("Messages sent to a model %", (0.0, 0.0, 0.0, 9.8))],
             colors=[BLACK, ORANGE], number_format='0.0"%"', max_scale=100)

chart_label(s, ML, 8.00, LW, "COST OF OWNERSHIP PER MILLION MESSAGES  —  AND WHERE IT ACTUALLY GOES")
table(s, ML, 8.28, 6.20, [
    ["COMPONENT", "PER MILLION", "NOTE"],
    ["Deterministic tiers 0–2 (CPU only)", "~$0.035", "94.7% of all decisions"],
    ["Model tier at ~5% escalation", "~$3.75", "$74/month at 10M messages"],
    [("Machine total", {"bold": True}), ("~$4", {"bold": True}), "Hosted, elastic, no ops"],
    [("Human review at 0.1%", {"bold": True}), ("~$125", {"bold": True, "color": ORANGE}),
     "30 s handling, $15/hr loaded"],
], col_w=[3.1, 1.3, 1.8], row_h=0.36, head_h=0.34, body_size=11)

doughnut_chart(s, ML + 6.50, 8.28, 4.90, 2.27,
               ["Machine — tiers 0–3", "Human review"], (4.0, 125.0),
               colors=[BLACK, ORANGE], number_format='"$"#,##0', label_size=10.5)

# right column
rx = ML + 11.95
card(s, rx, 2.70, 5.81, 2.38, "Humans cost ~30x the models",
     ["This inverts the usual instinct. The cheapest optimisation available is not a smaller model — "
      "it is better precision, because precision reduces review volume.",
      "Spending $4 of extra model budget to remove $40 of human review is correct. Minimising API "
      "calls gets the trade backwards."],
     bsize=12)

card(s, rx, 5.20, 5.81, 2.60, "Latency is spent where it is free",
     ["The user waits 3.45 ms. Nothing more.",
      "Tier 3's 1.7 s sits AFTER delivery, so it is invisible to the sender. Chat is already "
      "eventually consistent, so a verdict arriving a second later can still mask, withhold or retract.",
      "Worst case on deliberately pathological input is bounded and recorded in the verdict: "
      "14,000 chars → 114 ms; 29,400 chars → 168 ms."],
     bsize=12)

card(s, rx, 7.92, 5.81, 2.63, "Trade-offs we accepted deliberately",
     ["Redaction instead of blocking — cuts false-positive cost ~10x without touching detector accuracy.",
      "Lexicons deliberately incomplete — off-list means ASK, not ALLOW.",
      "We removed a working detector. Word-length counting decoded \u201cBreakfast is 250 per person, "
      "dinner around 600 for 2\u201d into a valid mobile number. It caused real false positives and "
      "caught only synthetic attacks, so it was deleted."],
     bsize=12)

footer(s, 7)

# ============================================================================
# SLIDE 8 — Safety guardrails
# ============================================================================
s = new_slide()
header(s, "07.", "Safety guardrails",
       "Contact protection is a revenue problem. Safety is a people problem. They share a pipeline and deliberately do NOT share failure semantics.")

# the asymmetry
box(s, ML, 2.70, 8.6, 1.62, fill=WHITE, line=ORANGE, lw=1.25)
text(s, ML + 0.3, 2.86, 8.0, 0.3, "THE GOVERNING ASYMMETRY", size=12, bold=True, color=ORANGE)
text(s, ML + 0.3, 3.24, 3.9, 0.9,
     [("CONTACT   uncertain → DELIVER", {"bold": True, "size": 13}),
      ("One lost commission is survivable. A stalled chat is not.",
       {"size": 11.5, "color": GREY})], line_spacing=1.05)
text(s, ML + 4.5, 3.24, 4.0, 0.9,
     [("SAFETY   uncertain → HOLD", {"bold": True, "size": 13, "color": ORANGE}),
      ("A delivered threat is not survivable. Held for human review.",
       {"size": 11.5, "color": GREY})], line_spacing=1.05)

card(s, ML, 4.50, 8.6, 1.85, "The target rule — the decision that makes this work",
     ['"this place is shit"  →  a crude review.  ALLOWED.',
      '"you are shit"  →  abuse.  WARNED.',
      "Same word, opposite verdict. The target decides, not the vocabulary. A guest describing a "
      "bad stay in strong language is exercising a right. Applied to phrase lists AND retrieval, "
      "because an insult aimed at an object shares nearly all its phrasing with one aimed at a person."],
     bsize=12)

card(s, ML, 6.53, 8.6, 2.00, "Enforcement ladder — blocking is the last resort",
     ["allow — deliver as sent, including angry complaints",
      "mask — deliver with offending spans redacted (default for contact sharing)",
      "warn — not sent, sender may edit and resend (harassment; often ends it, no moderator time)",
      "block — not sent, logged, actor risk incremented (threats, sexual, brand phishing)",
      "review — held for a human (blackmail, ambiguous scams, ALL safety abstentions)"],
     bsize=12)

card(s, ML, 8.71, 8.6, 1.87, "Safety overrides do not negotiate with thresholds",
     ["threat → block     sexual → block     coercion → review     harassment → warn",
      "scam ≥ 0.95 → block (brand impersonation is unambiguous)     scam < 0.95 → review",
      ("self-harm → DELIVER, and surface support resources. NEVER blocked — blocking a person "
       "in crisis is the worst available product behaviour.", {"color": ORANGE, "bold": True})],
     bsize=12)

# right column
rx = ML + 8.96
card(s, rx, 2.70, 8.80, 2.16, "Protecting the moderator itself  —  12/12 red-team caught",
     ["Messages are fenced and entity-escaped, so a sender cannot close the block and start "
      "issuing instructions. Role markers are defused; length is capped.",
      "Manipulation is enforced DETERMINISTICALLY — the model is never consulted for permission, "
      "because the component that might disagree is the one under attack.",
      "A message that tried to instruct the moderator can never be cleared BY the moderator."],
     bsize=12)

card(s, rx, 5.04, 8.80, 1.90, "The model's authority is deliberately limited",
     ["Says benign, but Tier 1 holds a validated number → IGNORED. It cannot overrule a fact.",
      "Says benign, but the message targeted the moderator → REFUSED, logged.",
      "Flags a coincidence-prone channel → needs high confidence. The same borderline input has "
      "been answered both ways across runs; enforcement must not ride on sampling variance.",
      "Returns two contradictory things → REJECTED → abstain → fails closed."],
     bsize=12)

card(s, rx, 7.12, 8.80, 1.72, "Availability guardrails",
     ["Breaker opens after 3 failures → 30 s cooldown → ONE half-open probe → success closes it; "
      "a failed probe doubles the cooldown, capped at 300 s.",
      "Per-minute AND per-day token buckets, self-resetting. Never a process-lifetime counter — "
      "a long-running worker would silently stop moderating forever."],
     bsize=12)

card(s, rx, 9.02, 8.80, 1.56, "Language coverage — verified, not claimed",
     ["English · romanised Hindi / Hinglish · Devanagari · Cyrillic. 18/18 unseen cases, zero "
      "false positives on difficult innocent controls in all four registers.",
      "Native-script pronouns match the ORIGINAL text: the homoglyph fold that defeats Cyrillic "
      "digit tricks also rewrites Cyrillic words, so matching the normalised view would work for "
      "Hindi and fail silently for Russian."],
     bsize=12)

footer(s, 8)

# ============================================================================
# SLIDE 9 — Proof, scale, roadmap
# ============================================================================
s = new_slide()
header(s, "08.", "Proof, scale and what comes next",
       "Every figure below is from a test run against the delivered code. None is a projection.")

LW = 11.40
chart_label(s, ML, 2.70, LW, "SUITE RESULTS")
table(s, ML, 3.00, LW, [
    ["SUITE", "RESULT"],
    [("Wayzyy official test cases", {"bold": True}), ("10 / 10", {"bold": True})],
    ["Regression suite — 103 cases (62 positive, 41 negative)", "100% recall · 100% precision · 0.00% FPR"],
    ["Adversarial suite — evasion levels L0–L5", "95.5% catch rate"],
    ["Red-team wave 2 — 119 novel attacks, 35 innocents",
     "61.3% caught by tiers 1+2 alone · 91.6% ceiling · 0/35 false positives"],
    ["Prompt-injection family", "12 / 12"],
    [("Unseen multilingual safety — 18 cases, 4 registers", {"bold": True}),
     ("18 / 18, zero false positives on innocent controls", {"bold": True})],
    ["Forced model outage (fault injection)", "Held for review — fail-closed verified"],
], col_w=[6.4, 5.0], row_h=0.38, head_h=0.34, body_size=11.5)

chart_label(s, ML, 6.12, 5.50,
            "WHICH TIER DECIDES  —  SHARE OF ALL VERDICTS.  BOTH ARE FREE.")
doughnut_chart(s, ML, 6.40, 5.50, 2.18,
               ["Tier 1 alone", "Tier 2 decides"], (94.7, 5.3),
               colors=[BLACK, ORANGE], number_format='0.0"%"', label_size=10.5)

chart_label(s, ML + 5.90, 6.12, 5.50, "MODEL-TIER COST PER MONTH BY VOLUME")
column_chart(s, ML + 5.90, 6.40, 5.50, 2.18,
             ["1M msgs", "10M msgs", "100M msgs"],
             [("Model tier cost per month", (7.40, 74.0, 740.0))],
             colors=[ORANGE], number_format='"$"#,##0', legend=False, label_size=10.0)

chart_label(s, ML, 8.78, LW, "SCALE  —  MEASURED ON ONE 8-CORE INSTANCE")
scale_metrics = [
    ("285", "msg/s per\nCPU core"),
    ("2,282", "msg/s on one\n8-core box"),
    ("~45,000", "concurrent active\nsenders / instance"),
    ("3.45 ms", "p50 write-path\nlatency"),
    ("197M", "messages/day\ntheoretical ceiling"),
]
for i, (v, lab) in enumerate(scale_metrics):
    metric(s, ML + i * 2.29, 9.04, 2.2, v, lab, vsize=23)

text(s, ML, 10.08, LW, 0.52,
     [("Scaling is horizontal and boring.", {"bold": True, "size": 11.5, "space": 1}),
      ("Tiers 0–2 are stateless — add instances, linear. Tier 3 is a durable queue with idempotent jobs. "
       "Only the conversation window and actor counters are shared state, both in Redis.",
       {"size": 11, "color": GREY, "space": 1})], line_spacing=1.02)

# right column — roadmap and honesty
rx = ML + 11.95
chart_label(s, rx, 2.70, 5.81, "ROADMAP")
roadmap = [
    ("1", "Labelled production evaluation set",
     "A golden set plus a difficult-innocent suite, held out of all tuning. Converts recall from a "
     "diagnostic into a reportable, drift-tracked metric."),
    ("2", "Resident multilingual safety classifier",
     "A quantised multi-label model on 100% of traffic at ~8 ms CPU — about 8 cents per million. "
     "Makes coverage language-general and permanently cuts model-tier volume."),
    ("3", "Provisional holds for critical severity",
     "Closes the ~1.7 s window in which a novel threat is visible before retraction. A product "
     "decision for Wayzyy, since it trades latency on a narrow slice of traffic."),
    ("4", "Actor and behavioural signals",
     "Velocity, repeat-target, fan-out, and treating user reports as a detection input rather than "
     "only a ticket. Plus cross-message aggregation: three sub-threshold messages become one case."),
]
ry = 3.02
for num, title, body in roadmap:
    chip(s, rx, ry, 0.36, 0.36, num, fill=ORANGE, size=13)
    text(s, rx + 0.52, ry + 0.02, 5.29, 0.30, title, size=12.5, bold=True, color=BLACK)
    text(s, rx + 0.52, ry + 0.34, 5.29, 0.76, body, size=11, color=GREY, line_spacing=1.0)
    ry += 1.22

box(s, rx, 7.96, 5.81, 0.035, fill=ORANGE)
text(s, rx, 8.14, 5.81, 1.20,
     [("What we are not claiming", {"bold": True, "size": 12.5, "color": ORANGE}),
      ("Recall on real Wayzyy traffic. The figures above are measured on adversarial corpora and "
       "targeted diagnostics, which demonstrate the mechanisms work but cannot substitute for a "
       "labelled production dataset. Building that set is item 1, and we expect to be held to the "
       "numbers it produces.", {"size": 11.5, "color": GREY})], line_spacing=1.05)

box(s, rx, 9.50, 5.81, 1.08, fill=BLACK)
text(s, rx + 0.3, 9.66, 5.2, 0.80,
     [("Reproduce every figure", {"bold": True, "size": 12, "color": WHITE}),
      ("./verify.sh · ./.verify/safety-audit · ./.verify/stats/tierstats",
       {"size": 11, "color": CREAM}),
      ("verify.sh fails the build if FPR rises above zero — enforced invariants, not metrics.",
       {"size": 10.5, "color": CREAM})], line_spacing=1.02, min_size=10.0)

footer(s, 9)

prs.save(OUT)
print(f"saved: {OUT}")
print(f"slides: {len(prs.slides.__iter__.__self__._sldIdLst)}")
